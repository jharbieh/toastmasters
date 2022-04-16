# Ref. https://python-pptx.readthedocs.io/en/latest/api/presentation.html
from pptx import Presentation     
import csv

# Create the pptx object
ppt = Presentation("d106-digital-banner.pptx")

# Get root slide
rslide = ppt.slides[0]
rslide_id = '{}'.format(rslide.slide_id)

print ("Got root slide: " + rslide_id)

# Selecting blank slide
blank_slide_layout = ppt.slide_layouts[2]

with open('clubs.csv', newline='') as csvfile:
    reader = csv.DictReader(csvfile)
    for row in reader:
        _division = row['Division']
        _area = row['Area']
        _id = row['ID']
        _name = row['Name']
        _status = row['Status']
        _charter = row['Date']
        _location = row['Location']

        # Attaching slide to ppt
        slide = ppt.slides.add_slide(blank_slide_layout) 

        # Adding title
        slide.shapes.title.text =  '{}'.format (_name)

        # Adding footer
        slide.shapes.placeholders[10].text = "Created By District 106 (C) 2022 For Spring Conference 2022"

        # Adding club number
        slide.shapes.placeholders[11].text = "Club ID: " + '{}'.format (_id)

        # Adding charter date
        slide.shapes.placeholders[12].text = "Club Charter Date: " + '{}'.format (_charter)

        # Adding area number
        slide.shapes.placeholders[13].text = "Club Area: " + '{}'.format (_area)

        # Adding division
        slide.shapes.placeholders[14].text = "Club Division: " + '{}'.format (_division)

        # Adding location
        slide.shapes.placeholders[15].text = "Club Location: " + '{}'.format (_location)

# save file
ppt.save('py.pptx')
  
print("Done")

