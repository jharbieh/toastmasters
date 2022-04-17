# Ref. https://python-pptx.readthedocs.io/en/latest/api/presentation.html
from pptx import Presentation     
from csv import csv

# Create the pptx object
ppt = Presentation("d106-digital-banner.pptx")

# Get root slide
rslide = ppt.slides[0]
rslide_id = '{}'.format(rslide.slide_id)

print ("Got root slide: " + rslide_id)

# Selecting blank slide
blank_slide_layout = ppt.slide_layouts[2]

# Attaching slide to ppt
slide = ppt.slides.add_slide(blank_slide_layout) 

# Slide placeholders
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))

# Adding title
slide.shapes.title.text = "Leaders Club Toastmasters"

# Ref. https://python-pptx.readthedocs.io/en/latest/user/placeholders-using.html
# Ref. https://python-pptx.readthedocs.io/en/latest/api/slides.html

# Adding footer
slide.shapes.placeholders[10].text = "Created By District 106 (C) 2022 For Spring Conference 2022"

# Adding club number
slide.shapes.placeholders[11].text = "Club ID: " + '{}'.format ("5342839")

# Adding charter date
slide.shapes.placeholders[12].text = "Club Charter Date: " + '{}'.format ("3/22/2016")

# Adding area number
slide.shapes.placeholders[13].text = "Club Area: " + '{}'.format ("23")

# Adding division
slide.shapes.placeholders[14].text = "Club Division: " + '{}'.format ("C")

# save file
ppt.save('py.pptx')
  
print("Done")

